# -*- coding: utf-8 -*-
"""
Created on Tue Jun 18 11:44:20 2019

@author: LaurencT
"""
import graph
import reshape_for_graphs
from pptx import Presentation
import os
import datetime

def draw_graphs_export(probabilistic_dict, deterministic_dict, bridge_graph_dict, graph_dir):
    """Writes graphs to graph_dir based on the data it is given
           Inputs:
               probabilistic_dict - dict - keys are codes and values are dfs
                   of trial parameters and estimates
               deterministic_dict - dict - keys are codes and values are dfs
                   of scenario parameters and estimates
               graph_dir - str - the directory where graphs should be exported
           Writes:
               5 graphs for each code to graph_dir
    """
    for code in probabilistic_dict.keys():
        graph_data = probabilistic_dict[code]
        variable = 'lives_touched'
        file_name = code + '_probabilistic_' + variable
        graph.probability_histogram(graph_data, variable, graph_dir, file_name)
        variable = 'lives_improved'
        file_name = code+'_probabilistic_'+variable
        graph.probability_histogram(graph_data, variable, graph_dir, file_name)
    
    for code in deterministic_dict.keys():
        param_df = deterministic_dict[code]
        variable = 'lives_touched'
        base, graph_data = reshape_for_graphs.restructure_graph_data_deterministic(param_df, variable)
        file_name = code+'_deterministic_'+variable
        graph.tornado_matplotlib(graph_data, base, graph_dir, file_name, variable)
        variable = 'lives_improved'
        base, graph_data = reshape_for_graphs.restructure_graph_data_deterministic(param_df, variable)
        file_name = code+'_deterministic_'+variable
        graph.tornado_matplotlib(graph_data, base, graph_dir, file_name, variable)
    
    for code in bridge_graph_dict.keys():
        file_name = code+'_bridge_graph'
        graph_data = bridge_graph_dict[code]
        graph.bridge_plot(graph_data, graph_dir, file_name)

# Section 13:
def graphs_to_slides(project_id, slides_dir, graph_dir, template_name):
    """Imports the Monday Meeting ppt template and adds graphs for a given project 
       creates a new dir in slides_dir for a new project and writes ppt there
       Inputs:
           project_id - str - an id_code for a model that has been run
           slides_dir - str - a directory to write the formatted ppt
           graph_dir - str - a directory where the graphs are saved
           template_name - str - the file name for the MM template ppt
       Writes:
           creates a directory in slides_dir and uploads a ppt file there
    """
    template_location = os.path.join(graph_dir, template_name)
    prs = Presentation(template_location)
        
    main_slide = prs.slides[0]
    
    bridge_path = os.path.join(graph_dir,project_id + '_bridge_graph.png')
    bridge_box = main_slide.shapes[5]
    bridge_top = bridge_box.top-10000
    bridge_left = bridge_box.left-10000
    main_slide.shapes.add_picture(bridge_path, 
                                  top = bridge_top, 
                                  left = bridge_left)
    
    lt_appendix_slide = prs.slides[1]
    
    lt_tornado_path = os.path.join(graph_dir, 
                                   project_id + '_deterministic_lives_touched.png')
    lt_tornado_box = lt_appendix_slide.shapes[5]
    lt_tornado_top = lt_tornado_box.top
    lt_tornado_left = lt_tornado_box.left-10000
    lt_appendix_slide.shapes.add_picture(lt_tornado_path, 
                                         top = lt_tornado_top, 
                                         left = lt_tornado_left)
    
    lt_histogram_path = os.path.join(graph_dir, 
                                     project_id + '_probabilistic_lives_touched.png')
    lt_histogram_box = lt_appendix_slide.shapes[6]
    lt_histogram_top = lt_histogram_box.top
    lt_histogram_left = lt_histogram_box.left
    lt_appendix_slide.shapes.add_picture(lt_histogram_path, 
                                         top = lt_histogram_top, 
                                         left = lt_histogram_left)
    
    li_appendix_slide = prs.slides[2]
    
    li_tornado_path = os.path.join(graph_dir, 
                                   project_id + '_deterministic_lives_improved.png')
    li_tornado_box = li_appendix_slide.shapes[5]
    li_tornado_top = li_tornado_box.top
    li_tornado_left = li_tornado_box.left-10000
    li_appendix_slide.shapes.add_picture(li_tornado_path, 
                                         top = li_tornado_top, 
                                         left = li_tornado_left)
    
    li_histogram_path = os.path.join(graph_dir, 
                                     project_id + '_probabilistic_lives_improved.png')
    li_histogram_box = li_appendix_slide.shapes[6]
    li_histogram_top = li_histogram_box.top-5000
    li_histogram_left = li_histogram_box.left
    li_appendix_slide.shapes.add_picture(li_histogram_path, 
                                         top = li_histogram_top, 
                                         left = li_histogram_left)
    
    new_dir = os.path.join(slides_dir, project_id)
    
    try:
        os.mkdir(new_dir)
    except FileExistsError:
        pass
    
    upload_path = os.path.join(new_dir+'/', 
                               project_id + '_mm_impact_charts.pptx')
    prs.save(upload_path)
    

def create_all_slides(param_dict, slides_dir, graph_dir, template_name):
    """Updates slides for each model that has been run"""
    for project_id in param_dict.keys():
        graphs_to_slides(project_id, slides_dir, graph_dir, template_name)

# Section 14:    
def update_estimates_output(deterministic_dict, probabilistic_dict, estimates_output, param_user):
    """Updates estimates_output with new values for lives_touched and lives_improved 
       (including the 95% confidence interval upper and lower bounds)
       Inputs:
           deterministic_dict - a dictionary where keys are id_codes and values
               are dfs where one of the indexes is 'base' and columns are 
               'lives_touched' and 'lives_improved'
           probabilistic_dict - a dictionary where keys are id_codes and values
               are dfs where columns are 'lives_touched' and 'lives_improved'
           estimates_output - a df where indexes are id_codes 
           param_user - dict - keys are id_codes for only the models that have been run
       Returns:
           a df updated with estimates
    """
    for code in param_user.keys():
        determ_df = deterministic_dict[code]
        prob_df = probabilistic_dict[code]
        estimates_output.loc[code, 'lives_touched'] = determ_df.loc['base', 'lives_touched']
        estimates_output.loc[code, 'lives_improved'] = determ_df.loc['base', 'lives_improved']
        estimates_output.loc[code, 'lives_touched_025'] = prob_df['lives_touched'].quantile(0.025)
        estimates_output.loc[code, 'lives_touched_975'] = prob_df['lives_touched'].quantile(0.975)
        estimates_output.loc[code, 'lives_improved_025'] = prob_df['lives_improved'].quantile(0.025)
        estimates_output.loc[code, 'lives_improved_975'] = prob_df['lives_improved'].quantile(0.975)
    return estimates_output


def export_estimates(estimates_output, analysis_type, backup_dir, outputs_dir, estimates_csv_name):
    """Overwrite the output csv with the updated estimates
       Inputs:
           param_user_all - a df of parameters and estimates
           analysis_type - a dict - keys are analysis options and values are 
               parameters, requires key for 'overwrite_estimates'
           backup_dir - str - a file path where backups of parameters should be written
           outputs_dir - str - a file path where output estimates should be written
           estimates_csv_name - str - the name of the output csv
       Writes:
           csvs to in the backup_dir and outputs_dir
    """
    # TODO Make sure these draw in the previous estimates and only update the relevant
    # ones that have been run
    # Write back up csv
    time_now = datetime.datetime.now()
    data_str = time_now.strftime('%Y_%m_%d_%H_%M')
    back_up_path = os.path.join(backup_dir, 'LTLI_parameters_python_' + data_str + '.csv')
    estimates_output.to_csv(back_up_path)
    # Overwrite imported csv
    if analysis_type['overwrite_estimates']:
        output_path = os.path.join(outputs_dir, estimates_csv_name)
        estimates_output.to_csv(output_path)